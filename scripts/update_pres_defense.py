from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


SRC = Path(r"C:\Users\Никита Жуковский\OneDrive\Рабочий стол\Жуковский Н.Е предзащита.pptx")


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.2), Inches(0.9))
    tf = t.text_frame
    tf.text = title
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.color.rgb = RGBColor(12, 63, 120)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.55), Inches(1.0), Inches(12.0), Inches(0.6))
        sf = s.text_frame
        sf.text = subtitle
        sf.paragraphs[0].font.size = Pt(15)
        sf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)


def add_bullets(slide, left: float, top: float, width: float, height: float, title: str, items: list[str]) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(242, 247, 255)
    box.line.color.rgb = RGBColor(170, 196, 230)
    tf = box.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(16)
    p0.font.color.rgb = RGBColor(28, 74, 132)
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.size = Pt(13)
        p.space_before = Pt(3)


def add_node(slide, x: float, y: float, w: float, h: float, text: str, color: tuple[int, int, int]) -> None:
    node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    node.fill.solid()
    node.fill.fore_color.rgb = RGBColor(*color)
    node.line.color.rgb = RGBColor(60, 90, 120)
    tf = node.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = RGBColor(80, 80, 80)
    c.line.width = Pt(1.4)


def slide_stack(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Технологический стек системы")
    add_bullets(
        s,
        0.5,
        1.5,
        5.9,
        2.2,
        "Frontend",
        [
            "React + TypeScript + Vite",
            "Recharts для интерактивных графиков",
            "Локальное хранилище для сценариев и настроек",
        ],
    )
    add_bullets(
        s,
        6.6,
        1.5,
        5.9,
        2.2,
        "Backend и данные",
        [
            "Python FastAPI (services: equipment, metrics)",
            "PostgreSQL",
            "REST API через gateway",
        ],
    )
    add_bullets(
        s,
        0.5,
        4.0,
        5.9,
        2.2,
        "Инфраструктура",
        [
            "Docker Compose",
            "Nginx (frontend + gateway)",
            "Разделение сервисов по ответственности",
        ],
    )
    add_bullets(
        s,
        6.6,
        4.0,
        5.9,
        2.2,
        "Инженерные практики",
        [
            "Ролевая модель и права доступа",
            "Отчетность (Excel + печать графиков)",
            "Модульная архитектура для масштабирования",
        ],
    )


def slide_component_uml(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "UML Диаграмма компонентов", "Логическое взаимодействие модулей системы")
    add_node(s, 0.5, 1.6, 2.6, 0.85, "Operator Panel", (222, 239, 255))
    add_node(s, 0.5, 2.8, 2.6, 0.85, "Main Web UI", (222, 239, 255))
    add_node(s, 3.6, 1.6, 2.8, 0.85, "Auth & Permissions", (232, 246, 228))
    add_node(s, 3.6, 2.8, 2.8, 0.85, "Reports Module", (232, 246, 228))
    add_node(s, 6.9, 1.6, 2.9, 0.85, "Gateway API", (255, 236, 217))
    add_node(s, 6.9, 2.8, 2.9, 0.85, "Metrics Service", (255, 236, 217))
    add_node(s, 10.2, 2.2, 2.1, 0.85, "Equipment Service", (255, 236, 217))
    add_node(s, 6.9, 4.2, 5.4, 0.95, "PostgreSQL + LocalStorage cache", (244, 232, 255))
    add_arrow(s, 3.1, 2.0, 3.6, 2.0)
    add_arrow(s, 3.1, 3.2, 3.6, 3.2)
    add_arrow(s, 6.4, 2.0, 6.9, 2.0)
    add_arrow(s, 6.4, 3.2, 6.9, 3.2)
    add_arrow(s, 9.8, 2.0, 10.2, 2.6)
    add_arrow(s, 8.3, 3.7, 8.3, 4.2)
    add_arrow(s, 11.2, 3.05, 11.2, 4.2)


def slide_sequence_uml(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "UML Диаграмма последовательности", "Формирование отчета эффективности/OEE")
    participants = [
        ("Пользователь", 0.8),
        ("Frontend", 3.2),
        ("Gateway", 5.6),
        ("Metrics Service", 8.0),
        ("Excel Export", 10.4),
    ]
    for name, x in participants:
        add_node(s, x, 1.4, 1.8, 0.6, name, (235, 243, 255))
        s.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x + 0.9), Inches(2.0), Inches(x + 0.9), Inches(6.5)
        ).line.color.rgb = RGBColor(150, 150, 150)
    steps = [
        (1.0, 4.1, "Запрос периода"),
        (3.4, 6.5, "GET /equipment + /metrics"),
        (5.8, 8.9, "Агрегация статусов"),
        (8.2, 10.9, "Набор показателей"),
        (6.5, 10.9, "XLSX генерация"),
        (10.9, 3.4, "Файл отчета"),
    ]
    y = 2.3
    for x1, x2, text in steps:
        add_arrow(s, x1, y, x2, y)
        lbl = s.shapes.add_textbox(Inches(min(x1, x2) + 0.1), Inches(y - 0.16), Inches(2.8), Inches(0.3))
        lbl.text_frame.text = text
        lbl.text_frame.paragraphs[0].font.size = Pt(10)
        y += 0.7


def slide_as_is(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "AS-IS: текущее состояние процессов (до внедрения)")
    add_node(s, 0.7, 2.0, 2.3, 0.9, "Ручной сбор\nданных", (255, 236, 217))
    add_node(s, 3.4, 2.0, 2.5, 0.9, "Фрагментированные\nисточники", (255, 236, 217))
    add_node(s, 6.3, 2.0, 2.5, 0.9, "Excel/бумажные\nжурналы", (255, 236, 217))
    add_node(s, 9.2, 2.0, 2.8, 0.9, "Запаздывающий\nанализ", (255, 236, 217))
    add_arrow(s, 3.0, 2.45, 3.4, 2.45)
    add_arrow(s, 5.9, 2.45, 6.3, 2.45)
    add_arrow(s, 8.8, 2.45, 9.2, 2.45)
    add_bullets(
        s,
        0.8,
        3.5,
        11.4,
        2.5,
        "Ключевые ограничения AS-IS",
        [
            "Нет единого окна мониторинга по станкам и печам",
            "Потеря точности при ручной фиксации простоев/операций",
            "Высокая трудоемкость формирования отчетности",
            "Ограниченная прозрачность эффективности оборудования",
        ],
    )


def slide_to_be(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "TO-BE: целевая модель после внедрения MDC-системы")
    add_node(s, 0.7, 2.0, 2.5, 0.9, "Поток данных\nв реальном времени", (232, 246, 228))
    add_node(s, 3.7, 2.0, 2.8, 0.9, "Единая платформа\nмониторинга", (232, 246, 228))
    add_node(s, 6.9, 2.0, 2.6, 0.9, "Автоматизированные\nотчеты", (232, 246, 228))
    add_node(s, 9.9, 2.0, 2.4, 0.9, "Управление\nэффективностью", (232, 246, 228))
    add_arrow(s, 3.2, 2.45, 3.7, 2.45)
    add_arrow(s, 6.5, 2.45, 6.9, 2.45)
    add_arrow(s, 9.5, 2.45, 9.9, 2.45)
    add_bullets(
        s,
        0.8,
        3.5,
        11.4,
        2.5,
        "Ключевые эффекты TO-BE",
        [
            "Онлайн-контроль загрузки и простоев оборудования",
            "Оперативная аналитика по печам и станкам с ПУ",
            "Единый контур отчетности: эффективность и OEE",
            "Ролевая модель доступа и прозрачность действий персонала",
        ],
    )


def slide_benefits(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Преимущества внедрения системы")
    add_bullets(
        s,
        0.6,
        1.6,
        12.0,
        4.7,
        "Бизнес- и производственные преимущества",
        [
            "Снижение времени реакции на отклонения по статусам оборудования",
            "Прозрачный учет регламентированных и нерегламентированных простоев",
            "Быстрое формирование отчетов (время, эффективность, OEE) в один клик",
            "Единый интерфейс для операторов, мастеров, термистов и администраторов",
            "Масштабируемость за счет модульной сервисной архитектуры",
        ],
    )


def slide_implementation(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Программная реализация и ключевые особенности")
    add_bullets(
        s,
        0.6,
        1.5,
        5.9,
        4.8,
        "Ключевые реализации",
        [
            "Сегментная модель статусов станков и печей",
            "Автоматические переходы фаз работы/простоя",
            "Симуляция метрик для ручных станков в операторской панели",
            "Гибкая настройка циклов станков с ПУ в админ-модуле",
        ],
    )
    add_bullets(
        s,
        6.7,
        1.5,
        5.9,
        4.8,
        "Особенности решения",
        [
            "Ролевая система с разграничением навигации и прав отчетности",
            "Интерактивные графики (масштаб, Brush, печать)",
            "Экспорт в Excel для управленческой отчетности",
            "Контейнеризация (Docker Compose) и раздельные сервисы",
        ],
    )


def main() -> None:
    if not SRC.exists():
        raise FileNotFoundError(f"Не найден файл: {SRC}")
    backup = SRC.with_name(f"{SRC.stem}.backup.{datetime.now().strftime('%Y%m%d-%H%M%S')}.pptx")
    backup.write_bytes(SRC.read_bytes())

    prs = Presentation(str(SRC))
    slide_stack(prs)
    slide_component_uml(prs)
    slide_sequence_uml(prs)
    slide_as_is(prs)
    slide_to_be(prs)
    slide_benefits(prs)
    slide_implementation(prs)
    prs.save(str(SRC))
    print(f"Updated: {SRC}")
    print(f"Backup:  {backup}")


if __name__ == "__main__":
    main()
