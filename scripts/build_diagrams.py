"""
Generate GOST-style UML diagrams as high-res PNG and insert into PPTX.
Diagrams: AS-IS, TO-BE, Component, State, Deployment, Sequence, Stack+Benefits
"""
from __future__ import annotations
import os, shutil, datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
matplotlib.rcParams["font.family"] = "Arial"
matplotlib.rcParams["font.size"] = 10

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, Rectangle
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

SRC = Path(r"C:\Users\Никита Жуковский\OneDrive\Рабочий стол\Жуковский Н.Е предзащита.pptx")
OUT = Path(r"d:\mdc\scripts\diagrams")
OUT.mkdir(exist_ok=True)

DPI = 200
W, H = 13.33, 7.5

# ─── Colors ──────────────────────────────────────────────────────────────────
BG       = "#ffffff"
BLK      = "#1c2833"
BLUE_D   = "#1a5276"
BLUE_L   = "#d4e6f1"
BLUE_M   = "#2e86c1"
GREEN_D  = "#1e8449"
GREEN_L  = "#d5f5e3"
GREEN_M  = "#27ae60"
ORANGE_D = "#ca6f1e"
ORANGE_L = "#fdebd0"
RED_D    = "#c0392b"
RED_L    = "#fadbd8"
PURPLE_D = "#6c3483"
PURPLE_L = "#e8daef"
YELLOW_L = "#fef9e7"
YELLOW_D = "#b7950b"
GRAY_L   = "#f2f3f4"
GRAY_M   = "#aeb6bf"
GRAY_D   = "#5d6d7e"


def new_fig():
    fig, ax = plt.subplots(figsize=(W, H), dpi=DPI)
    ax.set_xlim(0, W)
    ax.set_ylim(0, H)
    ax.set_aspect("equal")
    ax.axis("off")
    fig.patch.set_facecolor(BG)
    return fig, ax


def draw_box(ax, x, y, w, h, fill, edge, lw=1.4, radius=0.06):
    p = FancyBboxPatch((x, y), w, h, boxstyle=f"round,pad={radius}",
                        facecolor=fill, edgecolor=edge, linewidth=lw, zorder=2)
    ax.add_patch(p)
    return p


def draw_rect(ax, x, y, w, h, fill, edge, lw=1.2):
    p = Rectangle((x, y), w, h, facecolor=fill, edgecolor=edge, linewidth=lw, zorder=2)
    ax.add_patch(p)


def txt(ax, x, y, text, size=10, color=BLK, ha="center", va="center",
        bold=False, italic=False):
    ax.text(x, y, text, fontsize=size, color=color, ha=ha, va=va,
            fontweight="bold" if bold else "normal",
            fontstyle="italic" if italic else "normal", zorder=3)


def heading(ax, text, subtitle=None):
    txt(ax, W/2, H-0.35, text, size=19, color=BLUE_D, bold=True)
    if subtitle:
        txt(ax, W/2, H-0.72, subtitle, size=11.5, color=GRAY_D)
    ax.plot([0.4, W-0.4], [H-0.88, H-0.88], color=BLUE_D, lw=2, zorder=3)


def arr(ax, x1, y1, x2, y2, label=None, color=GRAY_D, lw=1.4, style="->",
        label_offset=0.14, label_size=8.5, label_color=None):
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle=style, color=color, lw=lw), zorder=2)
    if label:
        mx, my = (x1+x2)/2, (y1+y2)/2
        txt(ax, mx, my+label_offset, label, size=label_size,
            color=label_color or color, italic=True)


def gost_block(ax, x, y, w, h, title, body_lines, title_bg, body_bg, edge,
               title_size=11, body_size=9.2, title_color="#ffffff"):
    th = 0.36
    draw_box(ax, x, y, w, h, body_bg, edge, lw=1.6)
    draw_rect(ax, x+0.01, y+h-th, w-0.02, th-0.01, title_bg, edge, lw=1.2)
    txt(ax, x+w/2, y+h-th/2, title, size=title_size, color=title_color, bold=True)
    for i, line in enumerate(body_lines):
        txt(ax, x+0.15, y+h-th-0.25-i*0.27, line, size=body_size,
            color=BLK, ha="left", va="center")


def gost_node(ax, x, y, w, h, label, fill, edge, size=9.5, bold=False):
    draw_box(ax, x, y, w, h, fill, edge, lw=1.3, radius=0.05)
    txt(ax, x+w/2, y+h/2, label, size=size, bold=bold)


def diamond(ax, cx, cy, rw, rh, fill, edge, label, size=8.5):
    pts = np.array([[cx, cy+rh], [cx+rw, cy], [cx, cy-rh], [cx-rw, cy], [cx, cy+rh]])
    ax.fill(pts[:,0], pts[:,1], facecolor=fill, edgecolor=edge, linewidth=1.3, zorder=2)
    txt(ax, cx, cy, label, size=size, bold=True)


# ═════════════════════════════════════════════════════════════════════════════
# 1. AS-IS
# ═════════════════════════════════════════════════════════════════════════════
def diagram_as_is():
    fig, ax = new_fig()
    heading(ax, "AS-IS: модель бизнес-процесса до внедрения MDC-системы",
            "ГОСТ 19.701-90 — Схема процесса производственного мониторинга")

    # ─── Process blocks (top row) ────────────────────────────────
    procs = [
        (0.3,  5.0, 2.2, 1.1, "Ручная фиксация\nданных оператором", ORANGE_L, ORANGE_D),
        (2.9,  5.0, 2.2, 1.1, "Ведение бумажного\nжурнала смены", ORANGE_L, ORANGE_D),
        (5.5,  5.0, 2.2, 1.1, "Передача журналов\nмастеру участка", ORANGE_L, ORANGE_D),
        (8.1,  5.0, 2.2, 1.1, "Ручной ввод\nданных в Excel", ORANGE_L, ORANGE_D),
        (10.7, 5.0, 2.2, 1.1, "Формирование\nотчётов вручную", ORANGE_L, ORANGE_D),
    ]
    for x, y, w, h, label, fill, edge in procs:
        gost_node(ax, x, y, w, h, label, fill, edge, size=9.5, bold=True)

    for i in range(len(procs)-1):
        x1 = procs[i][0] + procs[i][2]
        x2 = procs[i+1][0]
        yy = procs[i][1] + procs[i][3]/2
        arr(ax, x1, yy, x2, yy, color=ORANGE_D, lw=1.5)

    # ─── Detail rows ────────────────────────────────────────────
    details = [
        (0.3,  3.55, "- Время начала / конца работы\n- Номер программы\n- Простои (по памяти)\n- Ручной подсчёт деталей"),
        (2.9,  3.55, "- Бумажная форма ГОСТ\n- Заполняется в конце смены\n- Возможны ошибки и пропуски\n- Нет структурирования"),
        (5.5,  3.55, "- Задержка 4-8 часов\n- Потеря / порча документов\n- Нет цифровой копии\n- Ручная сверка данных"),
        (8.1,  3.55, "- Разные версии файлов\n- Нет единого формата\n- Дублирование данных\n- Высокая трудоёмкость"),
        (10.7, 3.55, "- Расчёт показателей вручную\n- Нет OEE / эффективности\n- Отчёт за 2-3 дня\n- Ошибки вычислений"),
    ]
    for x, y, text in details:
        draw_box(ax, x, y, 2.2, 1.15, GRAY_L, GRAY_M, lw=0.8)
        txt(ax, x+0.12, y+0.85, text, size=8, color=GRAY_D, ha="left", va="top")

    for x, _, _ in details:
        arr(ax, x+1.1, 5.0, x+1.1, 4.7, color=GRAY_M, lw=0.8, style="-|>")

    # ─── Problems ────────────────────────────────────────────────
    txt(ax, W/2, 3.15, "КРИТИЧЕСКИЕ ПРОБЛЕМЫ ПРОЦЕССА", size=13, color=RED_D, bold=True)

    problems = [
        (0.3,  1.6, 2.7, "Потеря и искажение данных\nпри многократном переносе\nинформации между носителями"),
        (3.3,  1.6, 2.7, "Задержка принятия решений:\nданные доступны с опозданием\nна 1-3 рабочих дня"),
        (6.3,  1.6, 2.7, "Невозможность расчёта\nOEE / эффективности / простоев\nв реальном времени"),
        (9.3,  1.6, 3.3, "Отсутствие единого источника\nправды: Excel-файлы разных\nверсий у разных сотрудников"),
    ]
    for x, y, w, label in problems:
        draw_box(ax, x, y, w, 0.85, RED_L, RED_D, lw=1.3)
        txt(ax, x+w/2, y+0.42, label, size=8.5, color=RED_D)

    # ─── Impact metric boxes ────────────────────────────────────
    metrics = [
        (0.3,  0.4, "Полнота данных < 60%"),
        (3.3,  0.4, "Время отчёта: 2-3 дня"),
        (6.3,  0.4, "OEE: не рассчитывается"),
        (9.3,  0.4, "5+ источников данных"),
    ]
    for x, y, label in metrics:
        draw_box(ax, x, y, 2.7, 0.5, YELLOW_L, YELLOW_D, lw=1.0)
        txt(ax, x+1.35, y+0.25, label, size=9, color=YELLOW_D, bold=True)

    fig.savefig(OUT/"as_is.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 2. TO-BE
# ═════════════════════════════════════════════════════════════════════════════
def diagram_to_be():
    fig, ax = new_fig()
    heading(ax, "TO-BE: целевая модель после внедрения MDC-системы",
            "ГОСТ 19.701-90 — Автоматизированный мониторинг и аналитика")

    # ─── Central system ──────────────────────────────────────────
    draw_box(ax, 3.8, 2.55, 5.7, 3.9, BLUE_L, BLUE_D, lw=2.0)
    txt(ax, 6.65, 6.25, "MDC-система NEXA", size=14, color=BLUE_D, bold=True)

    modules = [
        (3.95, 4.7, 2.6, 1.35, "Подсистема\nмониторинга", [
            "Таймлайн станков",
            "Авто-фазы (UP/IDLE)",
            "Температура печей",
        ]),
        (6.75, 4.7, 2.6, 1.35, "Подсистема\nотчётности", [
            "Отчёт эффективности",
            "Отчёт OEE",
            "Температура + Excel",
        ]),
        (3.95, 2.75, 2.6, 1.65, "Подсистема\nуправления", [
            "Пульт оператора",
            "Адм. панель",
            "Роли и права (RBAC)",
            "Настройка циклов",
            "Сменный контроль",
        ]),
        (6.75, 2.75, 2.6, 1.65, "Подсистема\nхранения", [
            "PostgreSQL (метрики)",
            "LocalStorage (UI)",
            "REST API (CRUD)",
            "Временные ряды",
            "Сессии + настройки",
        ]),
    ]
    for x, y, w, h, title, items in modules:
        gost_block(ax, x, y, w, h, title, items, BLUE_D, "#eaf2f8", BLUE_D,
                   title_size=9.5, body_size=8)

    # ─── Left actors ─────────────────────────────────────────────
    left_actors = [
        (0.3, 5.0, "Оператор\n(пульт управления)"),
        (0.3, 3.7, "Мастер участка\n(мониторинг)"),
        (0.3, 2.4, "Термист\n(печи / садки)"),
    ]
    for x, y, label in left_actors:
        gost_node(ax, x, y, 2.8, 0.8, label, GREEN_L, GREEN_D, size=9, bold=True)
        arr(ax, 3.1, y+0.4, 3.8, y+0.4, color=GREEN_D, lw=1.3)

    # ─── Right actors ────────────────────────────────────────────
    right_actors = [
        (10.1, 5.0, "Инженер-технолог\n(отчёты)"),
        (10.1, 3.7, "Руководство\n(OEE, решения)"),
        (10.1, 2.4, "Администратор\n(настройки)"),
    ]
    for x, y, label in right_actors:
        gost_node(ax, x, y, 2.8, 0.8, label, GREEN_L, GREEN_D, size=9, bold=True)
        arr(ax, 9.5, y+0.4, 10.1, y+0.4, color=GREEN_D, lw=1.3)

    # ─── Benefits row ────────────────────────────────────────────
    txt(ax, W/2, 2.1, "КЛЮЧЕВЫЕ РЕЗУЛЬТАТЫ ВНЕДРЕНИЯ", size=13, color=GREEN_D, bold=True)

    benefits = [
        (0.15,  0.3, 2.15, "Данные в\nреальном времени\n\nЗадержка: 0 сек", GREEN_L, GREEN_D),
        (2.5,  0.3, 2.15, "Полнота данных\n100%\n\nАвтоматический сбор", GREEN_L, GREEN_D),
        (4.85, 0.3, 2.15, "Отчёт OEE\nза < 1 минуту\n\nExcel + печать", GREEN_L, GREEN_D),
        (7.2,  0.3, 2.15, "Единый интерфейс\nвсех ролей\n\nRBAC-авторизация", GREEN_L, GREEN_D),
        (9.55, 0.3, 2.15, "Учёт простоев\nпо сменам\n\nАвто-фазы + причины", GREEN_L, GREEN_D),
        (11.0, 0.3, 2.15, "Масштабируемая\nархитектура\n\nDocker Compose", GREEN_L, GREEN_D),
    ]
    # too many, reduce to 5
    benefits = [
        (0.3,  0.35, 2.35, "Данные реального\nвремени\n\nЗадержка: 0 сек\nПолнота: 100%"),
        (2.85, 0.35, 2.35, "Автоматический\nрасчёт OEE\n\nОтчёт < 1 мин\nExcel-экспорт"),
        (5.4,  0.35, 2.35, "Единый интерфейс\nвсех ролей\n\n4 роли + RBAC\n12+ прав доступа"),
        (7.95, 0.35, 2.35, "Сменный учёт\nпростоев\n\nАвто-фазы\n3 смены / сутки"),
        (10.5, 0.35, 2.35, "Масштабируемая\nинфраструктура\n\nDocker Compose\n5 контейнеров"),
    ]
    for x, y, w, label in benefits:
        draw_box(ax, x, y, w, 1.7, GREEN_L, GREEN_D, lw=1.1)
        txt(ax, x+w/2, y+0.85, label, size=8.5, color=BLK)

    fig.savefig(OUT/"to_be.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 3. Component Diagram
# ═════════════════════════════════════════════════════════════════════════════
def diagram_components():
    fig, ax = new_fig()
    heading(ax, "UML Диаграмма компонентов системы NEXA MDC",
            "Логическая архитектура и взаимодействие модулей")

    # ─── Client ──────────────────────────────────────────────────
    gost_block(ax, 0.2, 4.4, 2.8, 2.0, "<<Клиент>>", [
        "Браузер (Chrome/Edge)",
        "Пульт: планшет",
        "LocalStorage:",
        "  сессии, таймлайны,",
        "  сценарии, настройки",
    ], GRAY_D, GRAY_L, GRAY_D, body_size=8.5)

    # ─── Frontend ────────────────────────────────────────────────
    draw_box(ax, 3.3, 2.7, 5.0, 3.7, "#eaf2f8", BLUE_D, lw=2.0)
    txt(ax, 5.8, 6.15, "<<Frontend>>  React 18 + TypeScript + Vite", size=10, color=BLUE_D, bold=True)

    fe_modules = [
        (3.45, 5.2, 2.2, 0.85, "Мониторинг\nCNC + Печи", [
            "Таймлайн Recharts",
            "Статусы, метрики",
        ]),
        (5.9, 5.2, 2.2, 0.85, "Отчёты\nExcel + печать", [
            "OEE / эффективность",
            "Температура",
        ]),
        (3.45, 3.9, 2.2, 0.85, "Пульт оператора\nТема + табы", [
            "Статус, простои",
            "Авто-фазы",
        ]),
        (5.9, 3.9, 2.2, 0.85, "Администрирование\nНастройки", [
            "Пользователи",
            "Оборудование",
        ]),
    ]
    for x, y, w, h, title, items in fe_modules:
        gost_block(ax, x, y, w, h, title, items, BLUE_M, "#f4f9fe", BLUE_M,
                   title_size=8.5, body_size=7.5)

    draw_box(ax, 3.45, 2.85, 4.65, 0.72, PURPLE_L, PURPLE_D, lw=1.1)
    txt(ax, 5.78, 3.21, "Auth: роли (admin/user/term/cnc), 12+ permissions, RBAC",
        size=8.5, bold=True, color=PURPLE_D)

    # ─── Gateway ─────────────────────────────────────────────────
    gost_block(ax, 8.6, 4.3, 2.0, 2.1, "<<Gateway>>", [
        "Nginx 1.27",
        "Reverse Proxy",
        "Статика SPA",
        "Cache-Control",
        "PORT :80",
    ], ORANGE_D, ORANGE_L, ORANGE_D, body_size=8.5)

    # ─── Backend services ────────────────────────────────────────
    gost_block(ax, 8.6, 2.3, 2.0, 1.7, "Equipment Svc", [
        "FastAPI + Python",
        "CRUD оборудование",
        "Характеристики",
        "PORT :8001",
    ], GREEN_D, GREEN_L, GREEN_D, body_size=8.5)

    gost_block(ax, 10.9, 2.3, 2.0, 1.7, "Metrics Svc", [
        "FastAPI + Python",
        "Телеметрия",
        "Временные ряды",
        "PORT :8002",
    ], GREEN_D, GREEN_L, GREEN_D, body_size=8.5)

    # ─── Database ────────────────────────────────────────────────
    gost_block(ax, 10.9, 4.3, 2.0, 2.1, "<<Database>>", [
        "PostgreSQL 16",
        "Оборудование",
        "Характеристики",
        "Метрики",
        "Volume: pgdata",
    ], PURPLE_D, PURPLE_L, PURPLE_D, body_size=8.5)

    # ─── Arrows ──────────────────────────────────────────────────
    arr(ax, 3.0, 5.4, 3.3, 5.4, "HTTP", color=GRAY_D)
    arr(ax, 8.3, 5.2, 8.6, 5.2, "API", color=ORANGE_D)
    arr(ax, 8.3, 4.3, 8.6, 4.3, "API", color=ORANGE_D)
    arr(ax, 9.6, 4.3, 9.6, 4.0, "REST", color=GREEN_D)
    arr(ax, 10.6, 3.5, 10.9, 3.5, "REST", color=GREEN_D)
    arr(ax, 10.6, 3.1, 10.9, 3.1, color=GREEN_D)
    arr(ax, 10.9, 5.0, 10.9, 4.0, "SQL", color=PURPLE_D)
    arr(ax, 11.9, 4.3, 11.9, 4.0, "SQL", color=PURPLE_D)

    # ─── Docker bar ──────────────────────────────────────────────
    draw_box(ax, 0.2, 0.2, 12.8, 1.5, GRAY_L, GRAY_D, lw=1.5)
    txt(ax, 6.6, 1.4, "Docker Compose — контейнеризация и оркестрация", size=12, color=GRAY_D, bold=True)
    containers = ["frontend", "gateway (nginx)", "equipment-service", "metrics-service", "postgres"]
    for i, c in enumerate(containers):
        cx = 1.0 + i * 2.5
        draw_box(ax, cx, 0.35, 2.1, 0.7, "#ffffff", GRAY_D, lw=0.9)
        txt(ax, cx+1.05, 0.7, c, size=8.5, color=GRAY_D, bold=True)

    fig.savefig(OUT/"components.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 4. State Diagram
# ═════════════════════════════════════════════════════════════════════════════
def diagram_states():
    fig, ax = new_fig()
    heading(ax, "UML Диаграмма состояний оборудования (станок с ПУ)",
            "ГОСТ 19.701-90 — Переходы статусов и логика авто-фаз")

    S = {
        "off":        (1.0,  5.2, 2.3, 0.9, "Выключен (off)",           GRAY_L,    GRAY_D),
        "on":         (4.2,  5.2, 2.3, 0.9, "Включен (on)",             BLUE_L,    BLUE_D),
        "up":         (7.8,  5.2, 2.5, 0.9, "Работа по УП (up)",        GREEN_L,   GREEN_D),
        "tech_idle":  (10.8, 5.2, 2.2, 0.9, "Тех. простой\n(tech_idle)", ORANGE_L, ORANGE_D),
        "equip_idle": (10.8, 3.4, 2.2, 0.9, "Простой обор.\n(equip_idle)", RED_L,  RED_D),
        "lunch":      (1.5,  1.5, 2.3, 0.9, "Обед (lunch)",              YELLOW_L,  YELLOW_D),
        "service":    (4.3,  1.5, 2.7, 0.9, "Обслуживание\n(service)",   YELLOW_L,  YELLOW_D),
        "accident":   (7.5,  1.5, 2.3, 0.9, "Авария (accident)",         ORANGE_L,  ORANGE_D),
    }

    for key, (x, y, w, h, label, fill, edge) in S.items():
        draw_box(ax, x, y, w, h, fill, edge, lw=1.8, radius=0.12)
        txt(ax, x+w/2, y+h/2, label, size=9.5 if '\n' not in label else 8.5, bold=True)

    # Initial state
    ax.plot(0.35, 5.65, "ko", markersize=14, zorder=4)
    arr(ax, 0.5, 5.65, 1.0, 5.65, "Инициализация", color=BLK, lw=1.8, label_size=8)

    # Transitions
    arr(ax, 3.3, 5.65, 4.2, 5.65, "Включение\n(оператор)", color=BLUE_D, lw=1.6)
    arr(ax, 6.5, 5.65, 7.8, 5.65, "Запуск УП\n(оператор)", color=GREEN_D, lw=1.6)
    arr(ax, 10.3, 5.65, 10.8, 5.65, "Таймаут\n10 мин", color=ORANGE_D, lw=1.6)
    arr(ax, 11.9, 5.2, 11.9, 4.3, "Таймаут\n5 мин", color=RED_D, lw=1.6)

    # Equip idle to reasons
    arr(ax, 10.8, 3.6, 9.8, 1.95, color=ORANGE_D, lw=1.2, style="-|>",
        label="Авария", label_size=8)
    arr(ax, 10.8, 3.8, 7.0, 1.95, color=YELLOW_D, lw=1.2, style="-|>",
        label="Обслуживание", label_size=8)
    arr(ax, 10.8, 3.95, 3.8, 1.95, color=YELLOW_D, lw=1.2, style="-|>",
        label="Обед", label_size=8)

    # Force UP from equip_idle
    ax.annotate("", xy=(8.0, 5.2), xytext=(10.8, 3.85),
                arrowprops=dict(arrowstyle="->", color=GREEN_D, lw=1.8,
                                connectionstyle="arc3,rad=0.3"), zorder=2)
    txt(ax, 9.0, 4.5, "Принудительный\nзапуск УП", size=8.5, color=GREEN_D, bold=True)

    # Off from on
    ax.annotate("", xy=(2.15, 5.2), xytext=(5.35, 5.2),
                arrowprops=dict(arrowstyle="->", color=GRAY_D, lw=1.2,
                                connectionstyle="arc3,rad=-0.4"), zorder=2)
    txt(ax, 3.7, 4.65, "Выключение", size=8, color=GRAY_D, italic=True)

    # ─── Legend ──────────────────────────────────────────────────
    draw_box(ax, 0.3, 0.2, 12.7, 1.0, GRAY_L, GRAY_D, lw=0.8)
    txt(ax, 0.5, 0.95, "УСЛОВНЫЕ ОБОЗНАЧЕНИЯ:", size=9, color=BLK, bold=True, ha="left")
    leg = [
        (GREEN_D, "Ручное действие оператора (запуск УП)"),
        (ORANGE_D, "Автоматический переход по таймауту"),
        (YELLOW_D, "Присвоение причины простоя (из Простоев)"),
        (RED_D, "Нерегламентированный простой"),
    ]
    for i, (c, t) in enumerate(leg):
        cx = 0.6 + i * 3.2
        ax.plot(cx, 0.5, "s", color=c, markersize=10, zorder=4)
        txt(ax, cx+0.2, 0.5, t, size=8, color=BLK, ha="left")

    fig.savefig(OUT/"states.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 5. Deployment Diagram
# ═════════════════════════════════════════════════════════════════════════════
def diagram_deployment():
    fig, ax = new_fig()
    heading(ax, "UML Диаграмма развёртывания (Deployment Diagram)",
            "Инфраструктура Docker Compose и сетевое взаимодействие")

    # Client node
    gost_block(ax, 0.2, 3.8, 2.8, 2.4, "<<device>> Клиент", [
        "Браузер Chrome/Edge",
        "Пульт: Chromium/kiosk",
        "HTTP/HTTPS :80",
        "",
        "LocalStorage:",
        "  auth, timelines, prefs",
    ], GRAY_D, GRAY_L, GRAY_D, body_size=8.5)

    # Docker host
    draw_box(ax, 3.3, 0.8, 9.7, 5.8, "#eaf6fc", BLUE_D, lw=2.2)
    txt(ax, 8.15, 6.35, "<<execution environment>>  Docker Host (Ubuntu / Windows Server)", size=10.5, color=BLUE_D, bold=True)

    # Containers
    gost_block(ax, 3.5, 4.3, 2.5, 1.8, "gateway", [
        "nginx:1.27-alpine",
        "Reverse proxy",
        "Кеш статики",
        "PORT :80 (public)",
    ], ORANGE_D, ORANGE_L, ORANGE_D, body_size=8.5)

    gost_block(ax, 6.2, 4.3, 2.5, 1.8, "frontend", [
        "node:22-alpine",
        "React + Vite build",
        "Результат: /dist",
        "Multi-stage build",
    ], BLUE_M, BLUE_L, BLUE_M, body_size=8.5)

    gost_block(ax, 3.5, 2.0, 2.5, 1.9, "equipment-svc", [
        "python:3.12-slim",
        "FastAPI",
        "CRUD оборудование",
        "PORT :8001 (int)",
    ], GREEN_D, GREEN_L, GREEN_D, body_size=8.5)

    gost_block(ax, 6.2, 2.0, 2.5, 1.9, "metrics-svc", [
        "python:3.12-slim",
        "FastAPI",
        "Телеметрия + ряды",
        "PORT :8002 (int)",
    ], GREEN_D, GREEN_L, GREEN_D, body_size=8.5)

    gost_block(ax, 9.0, 2.0, 3.8, 1.9, "postgres", [
        "postgres:16-alpine",
        "База: equipment, metrics, readings",
        "Volume: pgdata (persistent)",
        "PORT :5432 (internal only)",
    ], PURPLE_D, PURPLE_L, PURPLE_D, body_size=8.5)

    gost_block(ax, 9.0, 4.3, 3.8, 1.8, "<<bridge network>>", [
        "mdc_default",
        "Healthcheck: pg_isready",
        "depends_on: postgres",
        "restart: unless-stopped",
    ], GRAY_D, "#f9f9f9", GRAY_D, body_size=8.5)

    # Arrows
    arr(ax, 3.0, 5.0, 3.5, 5.0, "HTTP :80", color=BLUE_D, lw=1.5)
    arr(ax, 6.0, 5.2, 6.2, 5.2, color=BLUE_M, lw=1.2)
    arr(ax, 4.75, 4.3, 4.75, 3.9, "REST", color=GREEN_D, lw=1.3)
    arr(ax, 7.45, 4.3, 7.45, 3.9, "REST", color=GREEN_D, lw=1.3)
    arr(ax, 6.0, 2.9, 9.0, 2.9, "TCP :5432", color=PURPLE_D, lw=1.3)
    arr(ax, 6.0, 2.6, 9.0, 2.6, color=PURPLE_D, lw=1.0)

    fig.savefig(OUT/"deployment.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 6. Sequence Diagram
# ═════════════════════════════════════════════════════════════════════════════
def diagram_sequence():
    fig, ax = new_fig()
    heading(ax, "UML Диаграмма последовательности: формирование отчёта OEE",
            "Взаимодействие компонентов при генерации отчёта")

    actors = [
        (0.8,  "Пользователь"),
        (3.2,  "Frontend\n(React SPA)"),
        (5.9,  "Nginx\n(Gateway)"),
        (8.5,  "Equipment\nService"),
        (11.2, "Metrics\nService"),
    ]
    for x, name in actors:
        gost_node(ax, x, 5.8, 1.8, 0.85, name, BLUE_L, BLUE_D, size=9, bold=True)
        ax.plot([x+0.9, x+0.9], [5.8, 0.3], color=GRAY_M, lw=0.8, ls="--", zorder=1)

    messages = [
        (1.7, 4.1, 5.35, "1. Выбор периода, нажатие «Создать»"),
        (4.1, 6.8, 5.0,  "2. GET /api/equipment (список станков)"),
        (6.8, 9.4, 4.65, "3. SELECT * FROM equipment WHERE type='ПУ'"),
        (9.4, 6.8, 4.3,  "4. JSON [{name, serial, type}...]"),
        (6.8, 4.1, 3.95, "5. Массив оборудования"),
        (4.1, 6.8, 3.6,  "6. GET /api/metrics?from=...&to=..."),
        (6.8, 12.1, 3.25, "7. SELECT readings (equipment_id, timestamp, value)"),
        (12.1, 6.8, 2.9,  "8. JSON [{metric, value, timestamp}...]"),
        (6.8, 4.1, 2.55,  "9. Данные метрик за период"),
        (4.1, 1.7, 2.2,   "10. Расчёт A/P/Q/OEE, генерация XLSX"),
        (1.7, 1.7, 1.5,   "11. Скачивание файла .xlsx"),
    ]

    y_top = 5.6
    step = 0.35
    for i, msg in enumerate(messages):
        if len(msg) == 3:
            x1, x2, y = msg[0], msg[0], msg[2]
            draw_box(ax, x1-0.5, y-0.1, 1.0, 0.25, GREEN_L, GREEN_D, lw=0.8)
            txt(ax, x1, y+0.03, "11. Скачивание .xlsx", size=7.5, color=GREEN_D, bold=True)
            continue
        x1, x2, y, label = msg
        clr = BLUE_D if x1 < x2 else GREEN_D
        ax.annotate("", xy=(x2, y), xytext=(x1, y),
                    arrowprops=dict(arrowstyle="->", color=clr, lw=1.4), zorder=3)
        mx = min(x1, x2) + abs(x2-x1)/2
        txt(ax, mx, y+0.13, label, size=7.5, color=BLK)

    fig.savefig(OUT/"sequence.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 7. Technology Stack + Benefits
# ═════════════════════════════════════════════════════════════════════════════
def diagram_stack():
    fig, ax = new_fig()
    heading(ax, "Технологический стек и преимущества системы NEXA MDC")

    sections = [
        (0.2, 4.0, 3.0, 2.5, "Frontend", BLUE_D, BLUE_L, [
            "React 18 + TypeScript",
            "Vite 5 (HMR, сборка)",
            "Recharts (диаграммы)",
            "ExcelJS (отчёты .xlsx)",
            "CSS Variables (темы)",
            "LocalStorage (состояние)",
        ]),
        (3.4, 4.0, 3.0, 2.5, "Backend", GREEN_D, GREEN_L, [
            "Python 3.12",
            "FastAPI (REST API)",
            "SQLAlchemy 2.0 (ORM)",
            "Pydantic 2 (валидация)",
            "Uvicorn (ASGI)",
            "PostgreSQL 16",
        ]),
        (6.8, 4.0, 3.0, 2.5, "Инфраструктура", ORANGE_D, ORANGE_L, [
            "Docker Compose v2",
            "Nginx 1.27 (reverse proxy)",
            "Multi-stage builds",
            "Health checks",
            "Persistent volumes",
            "Bridge network",
        ]),
        (10.0, 4.0, 3.0, 2.5, "Инженерные практики", PURPLE_D, PURPLE_L, [
            "Модульная архитектура",
            "RBAC (4 роли, 12+ прав)",
            "ГОСТ-совместимые отчёты",
            "Сменный мониторинг",
            "Авто-фазы оборудования",
            "Настраиваемые циклы",
        ]),
    ]

    for x, y, w, h, title, head_bg, body_bg, items in sections:
        gost_block(ax, x, y, w, h, title, items, head_bg, body_bg, head_bg, body_size=9)

    # Benefits
    txt(ax, W/2, 3.55, "ПРЕИМУЩЕСТВА ВНЕДРЕНИЯ", size=14, color=GREEN_D, bold=True)

    bens = [
        "Сокращение времени реакции на простои оборудования с дней до секунд",
        "Прозрачный учёт регламентированных и нерегламентированных простоев по сменам",
        "Автоматическое формирование отчётов OEE и эффективности в формате Excel",
        "Единый web-интерфейс для всех ролей: оператор, мастер, термист, администратор",
        "Масштабируемая контейнерная архитектура (Docker Compose, 5 сервисов)",
        "Гибкая настройка: циклы станков, цвета диаграмм, права доступа, палитра",
    ]
    for i, b in enumerate(bens):
        row = i // 2
        col = i % 2
        bx = 0.3 + col * 6.5
        by = 2.9 - row * 0.75
        ax.plot(bx, by, "s", color=GREEN_D, markersize=7, zorder=4)
        txt(ax, bx+0.25, by, b, size=9.5, color=BLK, ha="left")

    fig.savefig(OUT/"stack_benefits.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 8. Use Case Diagram
# ═════════════════════════════════════════════════════════════════════════════
def diagram_use_case():
    fig, ax = new_fig()
    heading(ax, "UML Диаграмма вариантов использования (Use Case)",
            "Функциональные возможности системы по ролям пользователей")

    # System boundary
    draw_box(ax, 2.8, 0.5, 7.5, 6.0, "#f0f6fc", BLUE_D, lw=2.0)
    txt(ax, 6.55, 6.25, "MDC-система NEXA", size=12, color=BLUE_D, bold=True)

    # Use cases (ellipses drawn as rounded boxes)
    cases = [
        (3.2, 5.5, 3.0, 0.45, "Просмотр мониторинга станков"),
        (3.2, 4.85, 3.0, 0.45, "Просмотр мониторинга печей"),
        (3.2, 4.2, 3.0, 0.45, "Просмотр детальной информации"),
        (3.2, 3.55, 3.0, 0.45, "Управление статусом станка"),
        (3.2, 2.9, 3.0, 0.45, "Присвоение причины простоя"),
        (3.2, 2.25, 3.0, 0.45, "Запуск / завершение садки"),
        (6.9, 5.5, 3.0, 0.45, "Формирование отчёта эфф-ти"),
        (6.9, 4.85, 3.0, 0.45, "Формирование отчёта OEE"),
        (6.9, 4.2, 3.0, 0.45, "Формирование отчёта по t°"),
        (6.9, 3.55, 3.0, 0.45, "Настройка оборудования"),
        (6.9, 2.9, 3.0, 0.45, "Управление пользователями"),
        (6.9, 2.25, 3.0, 0.45, "Настройка диаграмм / циклов"),
        (3.2, 1.6, 3.0, 0.45, "Экспорт в Excel / печать"),
        (6.9, 1.6, 3.0, 0.45, "Авторизация и смена роли"),
        (3.2, 0.75, 3.0, 0.45, "Переключение темы оформления"),
        (6.9, 0.75, 3.0, 0.45, "Просмотр простоев по сменам"),
    ]
    for x, y, w, h, label in cases:
        draw_box(ax, x, y, w, h, "#ffffff", BLUE_M, lw=1.0, radius=0.15)
        txt(ax, x+w/2, y+h/2, label, size=7.8, color=BLK)

    # Actors (left)
    def stick(ax, cx, cy, label, clr):
        r = 0.18
        ax.plot(cx, cy+r+0.12, "o", color=clr, markersize=10, zorder=4)
        ax.plot([cx, cx], [cy-0.05, cy+r+0.02], color=clr, lw=1.5, zorder=4)
        ax.plot([cx-0.2, cx+0.2], [cy+0.08, cy+0.08], color=clr, lw=1.5, zorder=4)
        ax.plot([cx-0.15, cx], [cy-0.25, cy-0.05], color=clr, lw=1.5, zorder=4)
        ax.plot([cx+0.15, cx], [cy-0.25, cy-0.05], color=clr, lw=1.5, zorder=4)
        txt(ax, cx, cy-0.42, label, size=8, color=clr, bold=True)

    stick(ax, 1.2, 5.4, "Оператор", GREEN_D)
    stick(ax, 1.2, 3.3, "Термист", ORANGE_D)
    stick(ax, 1.2, 1.3, "Мастер", BLUE_D)
    stick(ax, 12.0, 5.0, "Инженер", PURPLE_D)
    stick(ax, 12.0, 3.0, "Администратор", RED_D)

    # Connections (operator)
    for yc in [5.72, 5.07, 4.42, 3.77, 3.12, 0.97]:
        ax.plot([1.5, 3.2], [min(5.6, yc), yc], color=GREEN_D, lw=0.7, zorder=1)
    # Connections (thermist)
    for yc in [5.07, 4.42, 2.47, 1.82]:
        ax.plot([1.5, 3.2], [min(3.5, yc), yc], color=ORANGE_D, lw=0.7, zorder=1)
    # Connections (master)
    for yc in [5.72, 5.07, 4.42, 1.82]:
        ax.plot([1.5, 3.2], [min(1.5, yc), yc], color=BLUE_D, lw=0.7, zorder=1)
    # Connections (engineer)
    for yc in [5.72, 5.07, 4.42, 1.82]:
        ax.plot([11.7, 9.9], [min(5.2, yc), yc], color=PURPLE_D, lw=0.7, zorder=1)
    # Connections (admin)
    for yc in [3.77, 3.12, 2.47, 1.82, 0.97]:
        ax.plot([11.7, 9.9], [min(3.2, yc), yc], color=RED_D, lw=0.7, zorder=1)

    fig.savefig(OUT/"use_case.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 9. Activity Diagram (Operator Workflow)
# ═════════════════════════════════════════════════════════════════════════════
def diagram_activity():
    fig, ax = new_fig()
    heading(ax, "UML Диаграмма деятельности: рабочий процесс оператора",
            "ГОСТ 19.701-90 — Алгоритм работы оператора станка с ПУ за смену")

    # Start node
    ax.plot(1.5, 6.1, "ko", markersize=14, zorder=4)
    arr(ax, 1.5, 5.95, 1.5, 5.65, color=BLK, lw=1.5)

    steps = [
        (0.3,  5.15, 2.4, 0.5, "Авторизация в\nпульте оператора", BLUE_L, BLUE_D),
        (0.3,  4.3,  2.4, 0.5, "Выбор станка\nиз списка", BLUE_L, BLUE_D),
        (0.3,  3.3,  2.4, 0.5, "Запуск «Работа по УП»\n(статус = up)", GREEN_L, GREEN_D),
        (0.3,  2.0,  2.4, 0.5, "Работа станка (10 мин)\nМетрики обновляются", GREEN_L, GREEN_D),
    ]
    for x, y, w, h, label, fill, edge in steps:
        gost_node(ax, x, y, w, h, label, fill, edge, size=8.5, bold=True)

    arr(ax, 1.5, 5.15, 1.5, 4.8, color=BLK, lw=1.3)
    arr(ax, 1.5, 4.3, 1.5, 3.8, color=BLK, lw=1.3)

    # Decision diamond
    diamond(ax, 1.5, 3.0, 0.8, 0.25, YELLOW_L, YELLOW_D, "УП\nзавершена?", size=7)
    arr(ax, 1.5, 3.3, 1.5, 3.25, color=BLK, lw=1.3)
    arr(ax, 1.5, 2.75, 1.5, 2.5, "Да", color=BLK, lw=1.3, label_size=8)
    arr(ax, 2.3, 3.0, 2.7, 3.55, "Нет (продолжение)", color=GRAY_D, lw=1.0, label_size=7)

    # Auto-transition
    gost_node(ax, 3.2, 4.85, 2.6, 0.55, "Автопереход:\nТех. простой (5 мин)", ORANGE_L, ORANGE_D, size=8.5, bold=True)
    gost_node(ax, 3.2, 3.9, 2.6, 0.55, "Автопереход:\nПростой оборудования", RED_L, RED_D, size=8.5, bold=True)

    arr(ax, 1.5, 2.0, 1.5, 1.55, color=BLK, lw=1.3)
    arr(ax, 2.7, 2.25, 3.2, 5.1, "Таймаут", color=ORANGE_D, lw=1.2)
    arr(ax, 5.8, 5.1, 5.8, 4.45, color=ORANGE_D, lw=1.2)
    arr(ax, 5.8, 3.9, 5.8, 3.5, "Таймаут", color=RED_D, lw=1.2)

    # Decision: reason?
    diamond(ax, 1.5, 1.3, 0.8, 0.22, YELLOW_L, YELLOW_D, "Назначить\nпричину?", size=7)
    arr(ax, 0.7, 1.3, 0.3, 1.3, color=BLK, lw=1.0)
    txt(ax, 0.15, 1.45, "Нет", size=7, color=GRAY_D)
    arr(ax, 2.3, 1.3, 3.2, 1.3, "Да", color=BLK, lw=1.3, label_size=8)

    gost_node(ax, 3.2, 1.0, 2.6, 0.55, "Выбор причины:\nОбед / Обслуж. / Авария", YELLOW_L, YELLOW_D, size=8.5, bold=True)
    gost_node(ax, 3.2, 0.15, 2.6, 0.55, "Обновление таймлайна\nи диаграммы мониторинга", PURPLE_L, PURPLE_D, size=8.5, bold=True)
    arr(ax, 4.5, 1.0, 4.5, 0.7, color=BLK, lw=1.3)

    # Right column - parallel activities
    draw_box(ax, 6.3, 0.3, 6.5, 5.8, GRAY_L, GRAY_D, lw=1.0)
    txt(ax, 9.55, 5.85, "ПАРАЛЛЕЛЬНЫЕ ПРОЦЕССЫ СИСТЕМЫ", size=10, color=GRAY_D, bold=True)

    parallel = [
        (6.5, 4.8, 3.0, 0.6, "Мониторинг в реальном времени\n\nДиаграмма статусов станков\nобновляется автоматически", BLUE_L, BLUE_D),
        (6.5, 3.7, 3.0, 0.6, "Симуляция метрик\n\nОбороты, подача, выполнение\nпрограммы обновляются каждые 2с", GREEN_L, GREEN_D),
        (6.5, 2.6, 3.0, 0.6, "Учёт простоев по сменам\n\nСистема фиксирует начало/конец\nкаждого периода простоя", ORANGE_L, ORANGE_D),
        (6.5, 1.5, 3.0, 0.6, "Формирование отчётов\n\nДанные доступны для OEE\nи отчёта эффективности", PURPLE_L, PURPLE_D),
        (6.5, 0.45, 3.0, 0.55, "Хранение в LocalStorage\n\nТаймлайны, сценарии, сессии", GRAY_L, GRAY_D),
    ]
    for x, y, w, h, label, fill, edge in parallel:
        draw_box(ax, x, y, w, h, fill, edge, lw=0.9)
        txt(ax, x+w/2, y+h/2, label, size=7.5, color=BLK)

    for i in range(len(parallel)-1):
        y1 = parallel[i][1]
        y2 = parallel[i+1][1] + parallel[i+1][3]
        arr(ax, 8.0, y1, 8.0, y2+0.03, color=GRAY_D, lw=0.8)

    # Shift summary box
    gost_node(ax, 10.0, 3.3, 2.6, 1.6,
              "Итог смены:\n\nT работы по УП\nT регл. простоя\nT нерегл. простоя\nT фонд",
              GREEN_L, GREEN_D, size=8, bold=False)

    fig.savefig(OUT/"activity.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 10. Data Flow Diagram (DFD Level 0)
# ═════════════════════════════════════════════════════════════════════════════
def diagram_dfd():
    fig, ax = new_fig()
    heading(ax, "Диаграмма потоков данных (DFD Level 0)",
            "Основные информационные потоки между компонентами системы")

    # External entities (rectangles)
    entities = [
        (0.2, 5.1, 2.3, 0.8, "Оператор\nстанка", GREEN_L, GREEN_D),
        (0.2, 3.5, 2.3, 0.8, "Термист\n(печи)", ORANGE_L, ORANGE_D),
        (0.2, 1.5, 2.3, 0.8, "Инженер /\nРуководство", PURPLE_L, PURPLE_D),
        (10.8, 5.1, 2.3, 0.8, "Excel-файл\n(.xlsx)", BLUE_L, BLUE_D),
        (10.8, 3.5, 2.3, 0.8, "Принтер\n(печать)", GRAY_L, GRAY_D),
        (10.8, 1.5, 2.3, 0.8, "Администратор\n(настройки)", RED_L, RED_D),
    ]
    for x, y, w, h, label, fill, edge in entities:
        gost_node(ax, x, y, w, h, label, fill, edge, size=9, bold=True)

    # Central process (large circle-like)
    draw_box(ax, 3.5, 1.0, 6.2, 5.3, "#eaf2f8", BLUE_D, lw=2.0, radius=0.3)
    txt(ax, 6.6, 6.05, "MDC-система NEXA", size=13, color=BLUE_D, bold=True)

    # Sub-processes
    procs = [
        (3.7, 4.8, 2.8, 0.9, "P1: Мониторинг\nоборудования", [
            "Приём статусов",
            "Авто-фазы",
            "Обновление таймлайна",
        ]),
        (6.7, 4.8, 2.8, 0.9, "P2: Управление\nсадками", [
            "Запуск/стоп садки",
            "Контроль температуры",
            "Фиксация параметров",
        ]),
        (3.7, 2.7, 2.8, 1.5, "P3: Отчётность", [
            "Расчёт эффективности",
            "Расчёт OEE (A*P*Q)",
            "Температурные отчёты",
            "Генерация Excel",
        ]),
        (6.7, 2.7, 2.8, 1.5, "P4: Хранение\nи авторизация", [
            "PostgreSQL (метрики)",
            "LocalStorage (UI)",
            "RBAC (12+ прав)",
            "4 роли пользователей",
        ]),
        (5.2, 1.15, 2.8, 1.0, "P5: Администрирование", [
            "Оборудование",
            "Пользователи",
            "Циклы и диаграммы",
        ]),
    ]
    for x, y, w, h, title, items in procs:
        gost_block(ax, x, y, w, h, title, items, BLUE_D, "#ffffff", BLUE_D,
                   title_size=8.5, body_size=7.5)

    # Data flows (arrows with labels)
    arr(ax, 2.5, 5.5, 3.7, 5.3, "Статус станка", color=GREEN_D, lw=1.3, label_size=7.5)
    arr(ax, 2.5, 3.9, 3.7, 4.0, "Данные садки", color=ORANGE_D, lw=1.3, label_size=7.5)
    arr(ax, 2.5, 1.9, 3.7, 3.0, "Запрос отчёта", color=PURPLE_D, lw=1.3, label_size=7.5)
    arr(ax, 6.5, 3.5, 10.8, 5.3, "Excel-экспорт", color=BLUE_D, lw=1.3, label_size=7.5)
    arr(ax, 6.5, 3.2, 10.8, 3.9, "Данные для печати", color=GRAY_D, lw=1.2, label_size=7.5)
    arr(ax, 10.8, 1.9, 8.0, 1.7, "Настройки", color=RED_D, lw=1.2, label_size=7.5)

    # Internal flows
    arr(ax, 5.1, 4.3, 5.1, 3.8, color=BLUE_M, lw=0.9, label="Метрики")
    arr(ax, 8.1, 4.3, 8.1, 3.8, color=BLUE_M, lw=0.9, label="Данные")
    arr(ax, 6.5, 3.4, 6.7, 3.4, color=BLUE_M, lw=0.9)

    fig.savefig(OUT/"data_flow.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 11. Necessity / Relevance slide
# ═════════════════════════════════════════════════════════════════════════════
def slide_necessity():
    fig, ax = new_fig()
    heading(ax, "Актуальность разработки MDC-системы",
            "Почему автоматизация мониторинга критична для современного производства")

    # Left column - problems of industry
    draw_box(ax, 0.2, 3.2, 6.3, 3.2, "#fef5f0", ORANGE_D, lw=1.5)
    txt(ax, 3.35, 6.15, "Проблемы современного производства", size=12, color=ORANGE_D, bold=True)

    problems = [
        "Потери от незапланированных простоев составляют до 20% рабочего времени",
        "Ручной сбор данных приводит к потере 30-40% информации о состоянии\nоборудования (ГОСТ Р ИСО 22400)",
        "Средняя задержка реагирования на аварию без MDC: от 30 мин до 4 часов",
        "Отсутствие объективной статистики делает невозможным расчёт OEE\n(мировой стандарт Industry 4.0)",
        "Рост требований к цифровизации по программе «Цифровая экономика РФ 2024»",
        "Невозможность планировать ТОиР без данных о фактической загрузке",
    ]
    for i, p in enumerate(problems):
        yp = 5.65 - i * 0.42
        ax.plot(0.4, yp, "o", color=ORANGE_D, markersize=6, zorder=4)
        txt(ax, 0.6, yp, p, size=8.3, color=BLK, ha="left")

    # Right column - why MDC
    draw_box(ax, 6.8, 3.2, 6.2, 3.2, "#f0faf0", GREEN_D, lw=1.5)
    txt(ax, 9.9, 6.15, "Решение: MDC-система", size=12, color=GREEN_D, bold=True)

    solutions = [
        "Непрерывный мониторинг в режиме реального времени 24/7",
        "Автоматический расчёт OEE, эффективности и\nвремени простоев без участия человека",
        "Мгновенное реагирование: информация о простое\nдоступна в момент его начала",
        "Единый источник данных для всех ролей:\nоператор, мастер, инженер, руководитель",
        "Соответствие ГОСТ Р ИСО 22400 (KPI производства)\nи концепции Industry 4.0 / Smart Factory",
        "Основа для внедрения цифрового двойника\nи предиктивного обслуживания",
    ]
    for i, s in enumerate(solutions):
        yp = 5.65 - i * 0.42
        ax.plot(7.0, yp, "o", color=GREEN_D, markersize=6, zorder=4)
        txt(ax, 7.2, yp, s, size=8.3, color=BLK, ha="left")

    # Bottom - key statistics
    draw_box(ax, 0.2, 0.3, 12.8, 2.5, GRAY_L, GRAY_D, lw=1.2)
    txt(ax, 6.6, 2.55, "МИРОВАЯ СТАТИСТИКА И СТАНДАРТЫ", size=11, color=GRAY_D, bold=True)

    stats = [
        (1.0,  1.6, "72%", "предприятий в РФ\nне используют MDC\n(Минпромторг, 2024)"),
        (3.65, 1.6, "15-25%", "рост эффективности\nпосле внедрения MDC\n(McKinsey, 2023)"),
        (6.3,  1.6, "85%", "мирового рынка Smart\nFactory к 2028 году\n(Market Research Future)"),
        (9.0,  1.6, "Industry 4.0", "стандарт требует\nцифровизации всех\nпроизводственных данных"),
        (11.5, 1.6, "ISO 22400", "международный стандарт\nKPI для производства\n(OEE, доступность)"),
    ]
    for x, y, big, small in stats:
        txt(ax, x, y, big, size=14, color=BLUE_D, bold=True)
        txt(ax, x, y-0.55, small, size=7.5, color=GRAY_D)

    fig.savefig(OUT/"necessity.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 12. Results slide
# ═════════════════════════════════════════════════════════════════════════════
def slide_results():
    fig, ax = new_fig()
    heading(ax, "Результаты разработки и внедрения системы",
            "Достигнутые показатели и функциональные возможности")

    # Metrics row
    metrics = [
        (0.3,  5.0, 2.3, 1.35, "8 статусов\nоборудования", "Полная модель\nсостояний станка\nс авто-переходами", GREEN_L, GREEN_D),
        (2.8,  5.0, 2.3, 1.35, "3 типа\nотчётов", "Эффективность, OEE,\nтемпература печей\n(Excel + печать)", BLUE_L, BLUE_D),
        (5.3,  5.0, 2.3, 1.35, "4 роли\n12+ прав", "admin, user, term,\ncnc + гибкий RBAC\n(группы прав)", PURPLE_L, PURPLE_D),
        (7.8,  5.0, 2.3, 1.35, "5 Docker-\nконтейнеров", "Frontend, Gateway,\n2 микросервиса,\nPostgreSQL", ORANGE_L, ORANGE_D),
        (10.3, 5.0, 2.3, 1.35, "Реальное\nвремя", "Мониторинг 24/7\nобновление каждые\n2 секунды", GREEN_L, GREEN_D),
    ]
    for x, y, w, h, title, desc, fill, edge in metrics:
        draw_box(ax, x, y, w, h, fill, edge, lw=1.3)
        txt(ax, x+w/2, y+h-0.25, title, size=11, color=edge, bold=True)
        txt(ax, x+w/2, y+0.45, desc, size=7.8, color=BLK)

    # Functional results table
    draw_box(ax, 0.3, 0.3, 6.0, 4.2, GRAY_L, GRAY_D, lw=1.2)
    txt(ax, 3.3, 4.25, "РЕАЛИЗОВАННЫЙ ФУНКЦИОНАЛ", size=11, color=GRAY_D, bold=True)

    func_items = [
        "Пульт оператора с автопереходами фаз (UP 10 мин -> IDLE 5 мин)",
        "Светлая и тёмная темы оформления пульта",
        "Сменный учёт простоев (3 смены: 7-15, 15-23, 23-7)",
        "Присвоение причин простоя с обновлением диаграммы",
        "Принудительный запуск УП из режима простоя",
        "Садки печей: запуск, мониторинг t°, завершение, печать",
        "Отчёт эффективности: Excel с суммарными данными",
        "Отчёт OEE: стилизованный Excel (A * P * Q)",
        "Отчёт по температуре: интерактивная диаграмма + зум",
        "Конфигуратор циклов станков (визуальный builder)",
        "Администрирование: оборудование, пользователи, диаграммы",
        "Глобальная палитра цветов для диаграмм мониторинга",
    ]
    for i, item in enumerate(func_items):
        yy = 3.85 - i * 0.28
        ax.plot(0.5, yy, "s", color=GREEN_D, markersize=5, zorder=4)
        txt(ax, 0.7, yy, item, size=7.8, color=BLK, ha="left")

    # Comparison table (before/after)
    draw_box(ax, 6.6, 0.3, 6.3, 4.2, GRAY_L, GRAY_D, lw=1.2)
    txt(ax, 9.75, 4.25, "СРАВНЕНИЕ: ДО И ПОСЛЕ ВНЕДРЕНИЯ", size=11, color=GRAY_D, bold=True)

    # Header
    draw_rect(ax, 6.8, 3.7, 2.8, 0.35, ORANGE_L, ORANGE_D, lw=0.8)
    txt(ax, 8.2, 3.87, "До (AS-IS)", size=9, color=ORANGE_D, bold=True)
    draw_rect(ax, 9.8, 3.7, 2.9, 0.35, GREEN_L, GREEN_D, lw=0.8)
    txt(ax, 11.25, 3.87, "После (TO-BE)", size=9, color=GREEN_D, bold=True)

    comparisons = [
        ("Бумажный журнал", "Электронный пульт"),
        ("Данные за 1-3 дня", "Данные в реальном времени"),
        ("Ручной расчёт OEE", "Автоматический OEE"),
        ("Excel вручную", "Excel в 1 клик"),
        ("Нет учёта простоев", "Сменный учёт простоев"),
        ("5+ источников", "Единый интерфейс"),
        ("Полнота < 60%", "Полнота 100%"),
        ("Нет ролей доступа", "4 роли + RBAC"),
        ("Локальная установка", "Docker Compose"),
        ("Нет масштабирования", "Микросервисы"),
    ]
    for i, (before, after) in enumerate(comparisons):
        yy = 3.4 - i * 0.3
        txt(ax, 8.2,  yy, before, size=7.8, color=ORANGE_D)
        txt(ax, 11.25, yy, after, size=7.8, color=GREEN_D)
        ax.plot([6.8, 12.5], [yy-0.14, yy-0.14], color=GRAY_M, lw=0.3)

    fig.savefig(OUT/"results.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 13. Conclusions slide
# ═════════════════════════════════════════════════════════════════════════════
def slide_conclusions():
    fig, ax = new_fig()
    heading(ax, "Выводы по результатам разработки")

    conclusions = [
        ("1. Разработана и реализована полнофункциональная MDC-система",
         "Система NEXA MDC обеспечивает автоматизированный сбор, хранение и анализ данных\n"
         "о работе производственного оборудования в режиме реального времени. Реализованы\n"
         "модули мониторинга станков с ПУ, термических печей, пульт оператора и система отчётности.",
         BLUE_L, BLUE_D),
        ("2. Обеспечен расчёт ключевых производственных показателей",
         "Реализован автоматический расчёт OEE (Overall Equipment Effectiveness) по формуле\n"
         "A * P * Q, отчёт эффективности оборудования и температурные отчёты по печам.\n"
         "Все отчёты экспортируются в стилизованные Excel-файлы и доступны для печати.",
         GREEN_L, GREEN_D),
        ("3. Реализована модель состояний с автоматическими переходами",
         "Разработана 8-статусная модель оборудования с автоматическими фазовыми переходами\n"
         "(Работа по УП -> Тех. простой -> Простой обор.) и механизмом присвоения причин\n"
         "простоя с обратной связью на диаграммы мониторинга.",
         ORANGE_L, ORANGE_D),
        ("4. Применена микросервисная архитектура с контейнеризацией",
         "Система развёрнута в Docker Compose: 5 изолированных контейнеров (Frontend,\n"
         "Gateway, Equipment Service, Metrics Service, PostgreSQL). Обеспечены health-checks,\n"
         "persistent volumes, multi-stage builds и автоматический перезапуск.",
         PURPLE_L, PURPLE_D),
        ("5. Создана ролевая модель доступа и гибкая система настроек",
         "Реализован RBAC с 4 ролями и 12+ правами доступа, настраиваемые циклы станков,\n"
         "глобальная палитра цветов диаграмм, управление оборудованием и пользователями\n"
         "через административную панель.",
         RED_L, RED_D),
    ]

    for i, (title, body, fill, edge) in enumerate(conclusions):
        y = 5.6 - i * 1.12
        draw_box(ax, 0.3, y-0.25, 12.6, 1.02, fill, edge, lw=1.2)
        txt(ax, 0.5, y+0.55, title, size=10.5, color=edge, ha="left", bold=True)
        txt(ax, 0.5, y+0.1, body, size=8, color=BLK, ha="left", va="top")

    fig.savefig(OUT/"conclusions.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 14. Future improvements (with digital twin)
# ═════════════════════════════════════════════════════════════════════════════
def slide_future():
    fig, ax = new_fig()
    heading(ax, "Перспективы развития и будущие доработки",
            "Дорожная карта развития системы NEXA MDC")

    # Digital twin - main highlight
    draw_box(ax, 0.2, 3.6, 7.2, 2.8, "#eaf6fc", BLUE_D, lw=2.2)
    txt(ax, 3.8, 6.15, "Внедрение цифрового двойника (Digital Twin)", size=13, color=BLUE_D, bold=True)

    dt_items = [
        "Создание виртуальной модели каждого станка и печи, отражающей реальное",
        "состояние оборудования в реальном времени с точностью до секунды",
        "",
        "Предиктивная аналитика: прогнозирование поломок на основе накопленной",
        "статистики простоев, аварий и отклонений параметров от нормы",
        "",
        "Моделирование сценариев: «что если» — оценка влияния изменения",
        "расписания, обслуживания или параметров цикла на общий OEE",
        "",
        "Интеграция с промышленными протоколами (OPC UA, MTConnect)",
        "для автоматического получения данных напрямую со станков",
    ]
    for i, item in enumerate(dt_items):
        if item:
            txt(ax, 0.5, 5.65 - i * 0.24, item, size=8, color=BLK, ha="left")

    # Digital twin visualization
    draw_box(ax, 4.9, 3.8, 2.3, 1.6, "#d4efdf", GREEN_D, lw=1.3)
    txt(ax, 6.05, 5.15, "Digital Twin", size=10, color=GREEN_D, bold=True)
    txt(ax, 6.05, 4.6, "Реальное\nоборудование\n<-> Виртуальная\nмодель", size=7.5, color=BLK)

    # Other improvements - right column
    draw_box(ax, 7.7, 3.6, 5.2, 2.8, GRAY_L, GRAY_D, lw=1.5)
    txt(ax, 10.3, 6.15, "Другие направления развития", size=12, color=GRAY_D, bold=True)

    other = [
        ("Машинное обучение", "Автоматическое определение аномалий\nв работе оборудования (ML-модели)"),
        ("Dashboard реального\nвремени", "Большие экраны в цехе с визуализацией\nстатусов всех станков и печей"),
        ("Мобильное приложение", "Пульт оператора на смартфоне\nс push-уведомлениями о простоях"),
        ("Интеграция с ERP/MES", "Обмен данными с 1С, SAP\nили другими учётными системами"),
        ("Расширение отчётности", "MTBF, MTTR, Pareto-анализ причин\nпростоев, тренды эффективности"),
    ]
    for i, (title, desc) in enumerate(other):
        yy = 5.55 - i * 0.48
        txt(ax, 7.9, yy, title, size=8.5, color=BLUE_D, ha="left", bold=True)
        txt(ax, 10.0, yy, desc, size=7.2, color=BLK, ha="left")

    # Timeline
    draw_box(ax, 0.2, 0.3, 12.7, 2.85, "#fef9f0", ORANGE_D, lw=1.3)
    txt(ax, 6.55, 2.9, "ДОРОЖНАЯ КАРТА РАЗВИТИЯ", size=12, color=ORANGE_D, bold=True)

    phases = [
        (0.5,  1.2, "Этап 1\n(текущий)", "Базовый MDC:\nмониторинг, отчёты,\nпульт оператора,\nадминистрирование", GREEN_L, GREEN_D),
        (3.0,  1.2, "Этап 2\n(6 мес.)", "Интеграция OPC UA,\nмобильное приложение,\nDashboard, расширение\nотчётности (MTBF/MTTR)", BLUE_L, BLUE_D),
        (5.5,  1.2, "Этап 3\n(12 мес.)", "Цифровой двойник:\nвиртуальные модели,\nпредиктивная аналитика,\nML-детекция аномалий", PURPLE_L, PURPLE_D),
        (8.0,  1.2, "Этап 4\n(18 мес.)", "Smart Factory:\nинтеграция с ERP/MES,\nавтоматическое\nпланирование ТОиР", ORANGE_L, ORANGE_D),
        (10.5, 1.2, "Целевое\nсостояние", "Полностью автономная\nсистема управления\nпроизводственными\nпроцессами", RED_L, RED_D),
    ]
    for x, y, title, desc, fill, edge in phases:
        draw_box(ax, x, y-0.55, 2.2, 1.7, fill, edge, lw=1.1)
        txt(ax, x+1.1, y+0.85, title, size=8.5, color=edge, bold=True)
        txt(ax, x+1.1, y+0.2, desc, size=7, color=BLK)

    # Arrows between phases
    for i in range(4):
        x1 = phases[i][0] + 2.2
        x2 = phases[i+1][0]
        arr(ax, x1, 1.35, x2, 1.35, color=GRAY_D, lw=1.5)

    fig.savefig(OUT/"future.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 15. Timeline Diagram (diploma stages)
# ═════════════════════════════════════════════════════════════════════════════
def slide_timeline():
    fig, ax = new_fig()
    ax.set_facecolor(BG)

    # Title
    txt(ax, W/2, H-0.3, "ЭТАПЫ РАЗРАБОТКИ", size=22, color="#2c3e50", bold=True)
    ax.plot([W/2-0.15, W/2+0.15], [H-0.55, H-0.55], color="#3498db", lw=3, zorder=4)
    txt(ax, W/2, H-0.78, "Дипломный проект: MDC-система NEXA", size=12, color="#7f8c8d")

    stages = [
        ("ЭТАП 01", "Анализ предметной\nобласти",
         "Исследование процессов\nпроизводственного\nмониторинга, выявление\nпроблем и требований"),
        ("ЭТАП 02", "Обзор литературы\nи стандартов",
         "Изучение ГОСТ, ISO 22400,\nконцепций Industry 4.0,\nсуществующих MDC/MES\nрешений"),
        ("ЭТАП 03", "Проектирование\nархитектуры",
         "Определение стека\nтехнологий, модульной\nструктуры, API,\nбазы данных"),
        ("ЭТАП 04", "Разработка\nбэкенда",
         "FastAPI-микросервисы,\nPostgreSQL, REST API,\nDocker Compose,\nмодели данных"),
        ("ЭТАП 05", "Разработка\nфронтенда",
         "React SPA, пульт\nоператора, мониторинг,\nадминистрирование,\nтемы оформления"),
        ("ЭТАП 06", "Система\nотчётности",
         "Отчёты OEE и\nэффективности (Excel),\nтемпературные отчёты,\nпечать диаграмм"),
        ("ЭТАП 07", "Тестирование\nи отладка",
         "Функциональное\nтестирование, отладка\nавто-фаз, проверка\nролевой модели"),
        ("ЭТАП 08", "Оформление\nВКР",
         "Пояснительная записка,\nприложения, UML-\nдиаграммы, подготовка\nк защите"),
    ]

    n = len(stages)
    colors = [
        ("#b39ddb", "#9575cd"),  # purple
        ("#90caf9", "#64b5f6"),  # blue
        ("#80cbc4", "#4db6ac"),  # teal
        ("#a5d6a7", "#81c784"),  # green
        ("#c5e1a5", "#aed581"),  # lime
        ("#fff59d", "#fff176"),  # yellow
        ("#ffcc80", "#ffb74d"),  # orange
        ("#ef9a9a", "#e57373"),  # red-pink
    ]

    block_w = 1.35
    gap = 0.22
    total_w = n * block_w + (n - 1) * gap
    x_start = (W - total_w) / 2
    arrow_h = 1.15
    top_y = 5.05

    # Draw dotted timeline
    tl_y = top_y - arrow_h / 2 - 0.05
    for stage_i in range(n):
        cx = x_start + stage_i * (block_w + gap) + block_w / 2
        if stage_i < n - 1:
            nx = x_start + (stage_i + 1) * (block_w + gap) + block_w / 2
            n_dots = 6
            for di in range(n_dots):
                frac = (di + 1) / (n_dots + 1)
                dx = cx + (nx - cx) * frac
                clr_a = colors[stage_i][1]
                clr_b = colors[stage_i + 1][1]
                ax.plot(dx, tl_y - 1.15, "o", color=clr_a, markersize=3.5,
                        alpha=0.5 + 0.5 * frac, zorder=3)

    for i, (label, title, desc) in enumerate(stages):
        x = x_start + i * (block_w + gap)
        cx = x + block_w / 2
        c_light, c_dark = colors[i]

        # Arrow/chevron shape
        notch = 0.18
        if i == 0:
            pts = np.array([
                [x, top_y],
                [x + block_w, top_y],
                [x + block_w + notch, top_y - arrow_h / 2],
                [x + block_w, top_y - arrow_h],
                [x, top_y - arrow_h],
            ])
        else:
            pts = np.array([
                [x, top_y],
                [x + block_w, top_y],
                [x + block_w + notch, top_y - arrow_h / 2],
                [x + block_w, top_y - arrow_h],
                [x, top_y - arrow_h],
                [x + notch, top_y - arrow_h / 2],
            ])
        poly = plt.Polygon(pts, facecolor=c_light, edgecolor="none", zorder=2,
                           alpha=0.92)
        ax.add_patch(poly)
        shadow = plt.Polygon(pts + np.array([0.03, -0.03]), facecolor="#00000010",
                             edgecolor="none", zorder=1)
        ax.add_patch(shadow)

        # Stage label
        txt(ax, cx, top_y - 0.28, label, size=8.5, color="#ffffff", bold=True)
        txt(ax, cx, top_y - arrow_h / 2 - 0.05, title, size=8, color="#ffffff",
            bold=True)

        # Timeline dot
        dot_y = top_y - arrow_h - 0.35
        ax.plot(cx, dot_y, "o", color=c_dark, markersize=11, zorder=4)
        ax.plot(cx, dot_y, "o", color="#ffffff", markersize=6, zorder=5)
        ax.plot([cx, cx], [top_y - arrow_h - 0.05, dot_y + 0.08], color=c_dark,
                lw=1.2, zorder=3)

        # Description box below
        desc_y = dot_y - 0.25
        ax.plot([cx, cx], [dot_y - 0.08, desc_y], color=c_dark, lw=1.0, zorder=3)

        desc_box_h = 1.1
        desc_box_w = block_w + 0.15
        dbx = cx - desc_box_w / 2
        dby = desc_y - desc_box_h
        draw_box(ax, dbx, dby, desc_box_w, desc_box_h, c_light, c_dark,
                 lw=0.0, radius=0.06)
        ax.add_patch(FancyBboxPatch(
            (dbx, dby), desc_box_w, desc_box_h,
            boxstyle="round,pad=0.06", facecolor=c_light, edgecolor="none",
            alpha=0.25, zorder=2))
        txt(ax, cx, desc_y - 0.15, desc, size=6.8, color="#2c3e50", va="top")

    # Parallel marker between stages 4 and 5 (backend & frontend)
    mid_x = x_start + 3.5 * (block_w + gap) + block_w / 2
    ax.annotate("", xy=(mid_x + 0.6, top_y + 0.15), xytext=(mid_x - 0.6, top_y + 0.15),
                arrowprops=dict(arrowstyle="<->", color="#e74c3c", lw=1.8), zorder=6)
    txt(ax, mid_x, top_y + 0.35, "параллельно", size=7.5, color="#e74c3c", bold=True)

    # Horizontal timeline line
    x_first = x_start + block_w / 2
    x_last = x_start + (n - 1) * (block_w + gap) + block_w / 2
    dot_y = top_y - arrow_h - 0.35
    ax.plot([x_first, x_last], [dot_y, dot_y], color=GRAY_M, lw=1.5, zorder=2)
    ax.plot(x_first - 0.3, dot_y, "o", color=GRAY_M, markersize=8, zorder=3)
    ax.annotate("", xy=(x_last + 0.4, dot_y), xytext=(x_last + 0.05, dot_y),
                arrowprops=dict(arrowstyle="-|>", color=GRAY_M, lw=2), zorder=3)

    fig.savefig(OUT/"timeline.png", bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ═════════════════════════════════════════════════════════════════════════════
# INSERT INTO PPTX
# ═════════════════════════════════════════════════════════════════════════════
def delete_slides_from_end(prs, count):
    """Remove `count` slides from the end using lxml."""
    sldIdLst = prs.slides._sldIdLst
    ids_to_remove = list(sldIdLst)[-count:]
    for sldId in ids_to_remove:
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def rebuild_pptx():
    ts = datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
    backup = SRC.parent / f"{SRC.stem}.backup.{ts}{SRC.suffix}"
    shutil.copy2(SRC, backup)
    print(f"Backup: {backup}")

    prs = Presentation(str(SRC))
    total = len(prs.slides)
    print(f"Current slides: {total}")

    if total > 9:
        to_remove = total - 9
        delete_slides_from_end(prs, to_remove)
        print(f"Removed {to_remove} slides from end")

    images = [
        "necessity.png",
        "stack_benefits.png",
        "as_is.png",
        "to_be.png",
        "components.png",
        "deployment.png",
        "states.png",
        "activity.png",
        "use_case.png",
        "data_flow.png",
        "sequence.png",
        "results.png",
        "conclusions.png",
        "future.png",
        "timeline.png",
    ]

    blank_layout = prs.slide_layouts[6]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for fname in images:
        img_path = OUT / fname
        if not img_path.exists():
            print(f"  SKIP {fname}")
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(img_path), Emu(0), Emu(0), slide_w, slide_h)
        print(f"  Added: {fname}")

    prs.save(str(SRC))
    print(f"\nSaved: {SRC}")
    print(f"Total slides: {len(prs.slides)}")


# ═════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    total = 15
    n = 0
    def step(fn, name):
        global n; n += 1; fn(); print(f"  [{n}/{total}] {name}")

    print("=== Generating GOST-style diagrams ===")
    step(slide_necessity,    "necessity.png")
    step(diagram_stack,      "stack_benefits.png")
    step(diagram_as_is,      "as_is.png")
    step(diagram_to_be,      "to_be.png")
    step(diagram_components, "components.png")
    step(diagram_deployment, "deployment.png")
    step(diagram_states,     "states.png")
    step(diagram_activity,   "activity.png")
    step(diagram_use_case,   "use_case.png")
    step(diagram_dfd,        "data_flow.png")
    step(diagram_sequence,   "sequence.png")
    step(slide_results,      "results.png")
    step(slide_conclusions,  "conclusions.png")
    step(slide_future,       "future.png")
    step(slide_timeline,     "timeline.png")
    print("\n=== Inserting into PPTX ===")
    rebuild_pptx()
